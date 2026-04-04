from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                color=RGBColor(34, 34, 34), align=PP_ALIGN.LEFT, font_name="Aptos"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color
    p.alignment = align
    return box


def add_node(slide, x, y, r=0.08, fill=RGBColor(90, 98, 110), line=RGBColor(90, 98, 110)):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x - r), Inches(y - r), Inches(2 * r), Inches(2 * r)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_edge(slide, x1, y1, x2, y2, color=RGBColor(140, 146, 156), width=2.0):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_bar(slide, left, top, label, value, highlight=False):
    add_textbox(slide, left, top, Inches(0.28), Inches(0.16), label, size=11, color=RGBColor(72, 78, 86))
    bar_left = left + Inches(0.28)
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        bar_left,
        top + Inches(0.02),
        Inches(value),
        Inches(0.1),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(40, 127, 94) if highlight else RGBColor(167, 174, 184)
    bar.line.fill.background()
    if highlight:
        add_textbox(slide, bar_left + Inches(value) + Inches(0.04), top - Inches(0.01), Inches(0.45), Inches(0.16),
                    "best", size=10.5, bold=True, color=RGBColor(40, 127, 94))


def draw_clique_bridge_graph(slide, nodes, bridge=None, highlight_nodes=None, highlight_edges=None,
                             node_r=0.07, edge_color=RGBColor(142, 148, 158),
                             node_color=RGBColor(95, 102, 112), accent_color=RGBColor(221, 132, 46),
                             edge_width=1.8):
    highlight_nodes = set(highlight_nodes or [])
    highlight_edges = {tuple(sorted(e)) for e in (highlight_edges or [])}
    edges = [(0, 1), (1, 2), (0, 2), (3, 4), (4, 5), (3, 5)]
    if bridge is not None:
        edges.append(bridge)
    for i, j in edges:
        key = tuple(sorted((i, j)))
        color = accent_color if key in highlight_edges else edge_color
        width = edge_width + 0.8 if key in highlight_edges else edge_width
        add_edge(slide, nodes[i][0], nodes[i][1], nodes[j][0], nodes[j][1], color=color, width=width)
    for idx, (x, y) in enumerate(nodes):
        fill = accent_color if idx in highlight_nodes else node_color
        add_node(slide, x, y, r=node_r, fill=fill, line=fill)


def add_pill(slide, left, top, width, text, fill):
    pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.3))
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    pill.line.fill.background()
    add_textbox(slide, left, top + Inches(0.05), width, Inches(0.16), text, size=10.5, bold=True,
                color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)


def add_stage_arrow(slide, left, top, width, text):
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, width, Inches(0.36))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(235, 239, 242)
    arrow.line.fill.background()
    add_textbox(slide, left + Inches(0.08), top + Inches(0.08), width - Inches(0.16), Inches(0.16),
                text, size=12, bold=True, color=RGBColor(96, 104, 114), align=PP_ALIGN.CENTER)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(248, 247, 243)

title_color = RGBColor(32, 43, 56)
text_color = RGBColor(62, 68, 78)
muted = RGBColor(107, 114, 124)
accent = RGBColor(221, 132, 46)
green = RGBColor(40, 127, 94)
blue = RGBColor(74, 116, 184)
olive = RGBColor(124, 145, 73)

add_textbox(
    slide, Inches(0.7), Inches(0.34), Inches(12.0), Inches(0.45),
    "From Ecological Networks to Interpretable Structure, Prediction, and Repair",
    size=23, bold=True, color=title_color, font_name="Aptos Display"
)

add_stage_arrow(slide, Inches(0.82), Inches(1.0), Inches(3.5), "1. Discover")
add_stage_arrow(slide, Inches(4.92), Inches(1.0), Inches(3.5), "2. Validate")
add_stage_arrow(slide, Inches(9.02), Inches(1.0), Inches(3.5), "3. Repair")

col_lefts = [0.82, 4.92, 9.02]
col_w = 3.5

# Stage 1
x = col_lefts[0]
add_textbox(slide, Inches(x), Inches(1.5), Inches(col_w), Inches(0.45),
            "Search with ecologists for the decomposition that captures the key ecological motif.",
            size=14, color=text_color, align=PP_ALIGN.CENTER)
band1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.1), Inches(1.98), Inches(3.3), Inches(1.15))
band1.fill.solid()
band1.fill.fore_color.rgb = RGBColor(250, 250, 248)
band1.line.fill.background()
nodes1 = [
    (x + 0.72, 2.58), (x + 1.1, 2.26), (x + 1.39, 2.72),
    (x + 2.11, 2.58), (x + 2.49, 2.26), (x + 2.78, 2.72),
]
draw_clique_bridge_graph(slide, nodes1, bridge=(2, 3), highlight_nodes={2, 3}, highlight_edges={(2, 3)},
                         node_r=0.08, edge_width=2.2)
add_textbox(slide, Inches(x), Inches(3.24), Inches(col_w), Inches(0.2),
            "Candidate decompositions", size=12.5, bold=True, color=muted, align=PP_ALIGN.CENTER)
add_pill(slide, Inches(x + 0.15), Inches(3.32), Inches(0.95), "clique", olive)
add_pill(slide, Inches(x + 1.28), Inches(3.32), Inches(0.95), "bridge", blue)
add_pill(slide, Inches(x + 2.41), Inches(3.32), Inches(0.95), "clique+bridge", accent)
add_textbox(slide, Inches(x + 0.15), Inches(3.88), Inches(1.15), Inches(0.22),
            "discrete search", size=12.5, bold=True, color=title_color, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(x + 1.42), Inches(3.9), Inches(0.26), Inches(0.16),
            "<->", size=14, bold=True, color=muted, align=PP_ALIGN.CENTER)
expert = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 1.9), Inches(3.78), Inches(1.35), Inches(0.42))
expert.fill.solid()
expert.fill.fore_color.rgb = RGBColor(236, 240, 243)
expert.line.color.rgb = RGBColor(214, 219, 225)
add_textbox(slide, Inches(x + 1.95), Inches(3.9), Inches(1.25), Inches(0.14),
            "ecologist feedback", size=11.5, bold=True, color=title_color, align=PP_ALIGN.CENTER)
box1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.48), Inches(col_w), Inches(0.88))
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(251, 243, 232)
box1.line.fill.background()
add_textbox(slide, Inches(x + 0.14), Inches(4.72), Inches(col_w - 0.28), Inches(0.34),
            "Identify the bridge motif that ecologists recognise as meaningful.",
            size=15, bold=True, color=accent, align=PP_ALIGN.CENTER)

# Stage 2
x = col_lefts[1]
add_textbox(slide, Inches(x), Inches(1.5), Inches(col_w), Inches(0.45),
            "Treat each decomposition as a hypothesis and test what it predicts.",
            size=14, color=text_color, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(x), Inches(1.98), Inches(col_w), Inches(0.2),
            "Same interpretation across multiple networks", size=12.5, bold=True, color=muted, align=PP_ALIGN.CENTER)
band2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.08), Inches(2.18), Inches(3.34), Inches(0.95))
band2.fill.solid()
band2.fill.fore_color.rgb = RGBColor(250, 250, 248)
band2.line.fill.background()

def mini_graph(px, py):
    coords = [
        (px + 0.06, py + 0.12), (px + 0.28, py - 0.02), (px + 0.36, py + 0.22),
        (px + 0.62, py + 0.12), (px + 0.84, py - 0.02), (px + 0.92, py + 0.22),
    ]
    draw_clique_bridge_graph(slide, coords, bridge=(2, 3), highlight_nodes={2, 3}, highlight_edges={(2, 3)},
                             node_r=0.05, edge_width=1.4)

mini_graph(x + 0.18, 2.42)
mini_graph(x + 1.25, 2.42)
mini_graph(x + 2.32, 2.42)
for i, lbl in enumerate(["instance 1", "instance 2", "instance 3"]):
    add_textbox(slide, Inches(x + 0.18 + i * 1.07), Inches(3.02), Inches(0.88), Inches(0.16),
                lbl, size=10.5, color=muted, align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(x + 0.12), Inches(3.42), Inches(1.15), Inches(0.2),
            "Outcome", size=12.5, bold=True, color=muted, align=PP_ALIGN.CENTER)
outcome = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.05), Inches(3.68), Inches(1.28), Inches(0.62))
outcome.fill.solid()
outcome.fill.fore_color.rgb = RGBColor(244, 246, 248)
outcome.line.color.rgb = RGBColor(220, 224, 229)
add_textbox(slide, Inches(x + 0.18), Inches(3.91), Inches(1.02), Inches(0.16),
            "dispersal rate", size=15, bold=True, color=title_color, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(x + 1.46), Inches(3.92), Inches(0.24), Inches(0.16),
            "->", size=15, bold=True, color=muted, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(x + 2.02), Inches(3.42), Inches(1.32), Inches(0.2),
            "Predictive score", size=12.5, bold=True, color=muted, align=PP_ALIGN.CENTER)
score = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 1.82), Inches(3.68), Inches(1.55), Inches(0.9))
score.fill.solid()
score.fill.fore_color.rgb = RGBColor(244, 246, 248)
score.line.color.rgb = RGBColor(220, 224, 229)
add_bar(slide, Inches(x + 2.0), Inches(3.88), "A", 0.6, False)
add_bar(slide, Inches(x + 2.0), Inches(4.13), "B", 1.0, True)
add_bar(slide, Inches(x + 2.0), Inches(4.38), "C", 0.76, False)
box2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.48), Inches(col_w), Inches(0.88))
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(235, 244, 239)
box2.line.fill.background()
add_textbox(slide, Inches(x + 0.14), Inches(4.72), Inches(col_w - 0.28), Inches(0.34),
            "The best predictive interpretation is the strongest structural hypothesis.",
            size=15, bold=True, color=green, align=PP_ALIGN.CENTER)

# Stage 3
x = col_lefts[2]
add_textbox(slide, Inches(x), Inches(1.5), Inches(col_w), Inches(0.45),
            "Use the winning structure to constrain completion, repair, and intervention.",
            size=14, color=text_color, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(x + 0.1), Inches(1.98), Inches(1.1), Inches(0.2),
            "Damaged network", size=12.5, bold=True, color=muted, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(x + 2.1), Inches(1.98), Inches(1.1), Inches(0.2),
            "Repaired network", size=12.5, bold=True, color=muted, align=PP_ALIGN.CENTER)
band3 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.05), Inches(2.18), Inches(3.4), Inches(1.0))
band3.fill.solid()
band3.fill.fore_color.rgb = RGBColor(250, 250, 248)
band3.line.fill.background()
left_nodes = [
    (x + 0.32, 2.62), (x + 0.7, 2.3), (x + 0.99, 2.76),
    (x + 1.3, 2.62), (x + 1.68, 2.3), (x + 1.97, 2.76),
]
right_nodes = [
    (x + 1.53, 2.62), (x + 1.91, 2.3), (x + 2.2, 2.76),
    (x + 2.82, 2.62), (x + 3.2, 2.3), (x + 3.49, 2.76),
]
draw_clique_bridge_graph(slide, left_nodes, bridge=None, node_r=0.075, edge_width=1.9)
add_textbox(slide, Inches(x + 1.18), Inches(2.56), Inches(0.24), Inches(0.16),
            "->", size=16, bold=True, color=muted, align=PP_ALIGN.CENTER)
draw_clique_bridge_graph(slide, right_nodes, bridge=(2, 3), highlight_nodes={2, 3}, highlight_edges={(2, 3)},
                         node_r=0.075, edge_width=1.9)
add_textbox(slide, Inches(x), Inches(3.42), Inches(col_w), Inches(0.34),
            "Constrain generation so the bridge motif is restored or preserved.",
            size=13.5, color=text_color, align=PP_ALIGN.CENTER)
box3 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.48), Inches(col_w), Inches(0.88))
box3.fill.solid()
box3.fill.fore_color.rgb = RGBColor(235, 244, 239)
box3.line.fill.background()
add_textbox(slide, Inches(x + 0.14), Inches(4.72), Inches(col_w - 0.28), Inches(0.34),
            "Use structure-guided generation to complete, repair, or improve ecological networks.",
            size=15, bold=True, color=green, align=PP_ALIGN.CENTER)

footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(6.72), Inches(11.7), Inches(0.44))
footer.fill.solid()
footer.fill.fore_color.rgb = RGBColor(233, 239, 235)
footer.line.fill.background()
add_textbox(
    slide, Inches(1.05), Inches(6.84), Inches(11.24), Inches(0.18),
    "Search with experts -> validate by prediction -> repair and design with structure",
    size=15, bold=True, color=title_color, align=PP_ALIGN.CENTER
)

out_path = "slides/ecological_networks_structure_prediction_repair.pptx"
prs.save(out_path)
print(out_path)
